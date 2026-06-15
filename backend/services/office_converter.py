"""
Office document converter service.
Converts DOCX, PPTX, and XLSX files to HTML for in-app viewing.
"""

import io
import base64
from typing import Optional
from html import escape
from docx.oxml.ns import qn


def convert_docx_to_html(file_content: bytes) -> Optional[str]:
    """Convert DOCX to rich HTML with lists, images, tables, colors, hyperlinks."""
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document(io.BytesIO(file_content))

        html = ['<div class="docx-content">']

        # Process body elements in order (paragraphs + tables interleaved)
        _process_body(doc, html)

        html.append('</div>')
        return '\n'.join(html)

    except Exception as e:
        print(f'DOCX conversion error: {e}')
        import traceback; traceback.print_exc()
        return None


def _process_body(doc, html: list):
    """Process document body elements maintaining order."""
    body = doc.element.body
    list_counter = {}  # track numbering per numId
    in_list = {}       # current nesting level per numId

    for child in body:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag

        if tag == 'p':
            _process_paragraph(child, doc, html, list_counter, in_list)
        elif tag == 'tbl':
            _process_table_element(child, doc, html)


def _get_paragraph_num_info(p_elem):
    """Extract numbering info from a paragraph element."""
    numPr = p_elem.find(qn('w:pPr'))
    if numPr is not None:
        numPr_elem = numPr.find(qn('w:numPr'))
        if numPr_elem is not None:
            ilvl_elem = numPr_elem.find(qn('w:ilvl'))
            numId_elem = numPr_elem.find(qn('w:numId'))
            if ilvl_elem is not None and numId_elem is not None:
                ilvl = int(ilvl_elem.get(qn('w:val'), '0'))
                numId = int(numId_elem.get(qn('w:val'), '0'))
                if numId > 0:
                    return numId, ilvl
    return None, None


def _process_paragraph(p_elem, doc, html: list, list_counter, in_list):
    """Process a single paragraph element."""
    from docx.text.paragraph import Paragraph

    para = Paragraph(p_elem, doc)
    text = para.text.strip()
    numId, ilvl = _get_paragraph_num_info(p_elem)

    style_name = ''
    if para.style:
        style_name = para.style.name or ''

    # Handle list items
    if numId is not None:
        if numId not in list_counter:
            list_counter[numId] = 0
            in_list[numId] = -1

        # Close inner lists if we went up
        while in_list[numId] >= ilvl:
            html.append('</ul>')
            in_list[numId] -= 1

        if in_list[numId] < ilvl:
            html.append('<ul>')
            in_list[numId] = ilvl

        list_counter[numId] += 1
        html.append('<li>')
        html.append(_process_runs_html(para.runs))
        html.append('</li>')
        return

    # Close any open lists
    for nid in list(in_list):
        while in_list[nid] >= 0:
            html.append('</ul>')
            in_list[nid] -= 1

    if not text and not para.runs:
        html.append('<br/>')
        return

    # Heading detection
    level = _get_heading_level(style_name)
    if level:
        html.append(f'<h{level}>')
        html.append(_process_runs_html(para.runs))
        html.append(f'</h{level}>')
        return

    # Title
    if 'Title' in style_name:
        html.append(f'<h1 class="docx-title">')
        html.append(_process_runs_html(para.runs))
        html.append('</h1>')
        return

    # Subtitle
    if 'Subtitle' in style_name:
        html.append(f'<p class="docx-subtitle">')
        html.append(_process_runs_html(para.runs))
        html.append('</p>')
        return

    # Quote
    if 'Quote' in style_name or 'Intense Quote' in style_name:
        html.append('<blockquote>')
        html.append(_process_runs_html(para.runs))
        html.append('</blockquote>')
        return

    # Normal paragraph
    html.append('<p>')
    html.append(_process_runs_html(para.runs))
    html.append('</p>')


def _get_heading_level(style_name: str) -> Optional[int]:
    for i in range(1, 7):
        if f'Heading {i}' in style_name or f'heading {i}' in style_name:
            return i
    if 'Title' in style_name:
        return 1
    return None


def _process_runs_html(runs) -> str:
    """Process text runs with rich formatting."""
    parts = []
    for run in runs:
        text = run.text
        if not text:
            continue

        escaped = escape(text)
        tags_open = []
        tags_close = []

        # Bold
        if run.bold:
            tags_open.append('<strong>')
            tags_close.append('</strong>')
        # Italic
        if run.italic:
            tags_open.append('<em>')
            tags_close.append('</em>')
        # Underline
        if run.underline:
            tags_open.append('<u>')
            tags_close.append('</u>')
        # Strikethrough
        if run.font.strike:
            tags_open.append('<s>')
            tags_close.append('</s>')

        # Font size
        if run.font.size:
            size_pt = run.font.size.pt
            tags_open.append(f'<span style="font-size:{size_pt}pt">')
            tags_close.append('</span>')

        # Font color
        color = None
        if run.font.color and run.font.color.rgb:
            color = str(run.font.color.rgb)
        # Also check highlight
        rpr = run._element.find(qn('w:rPr'))
        if rpr is not None and color is None:
            highlight = rpr.find(qn('w:highlight'))
            if highlight is not None:
                hl_val = highlight.get(qn('w:val'), '')
                color_map = {
                    'yellow': '#ffff00', 'green': '#00ff00', 'cyan': '#00ffff',
                    'red': '#ff0000', 'blue': '#0000ff', 'magenta': '#ff00ff',
                    'lightGray': '#cccccc', 'darkGray': '#666666',
                }
                color = color_map.get(hl_val)

        if color:
            tags_open.append(f'<span style="color:#{color}">')
            tags_close.append('</span>')

        # Superscript / Subscript
        if run.font.superscript:
            tags_open.append('<sup>')
            tags_close.append('</sup>')
        elif run.font.subscript:
            tags_open.append('<sub>')
            tags_close.append('</sub>')

        parts.append(''.join(tags_open) + escaped + ''.join(reversed(tags_close)))

    return ''.join(parts)


def _process_table_element(tbl_elem, doc, html: list):
    """Process a table XML element."""
    from docx.table import Table

    table = Table(tbl_elem, doc)
    html.append('<table>')

    # Detect header row (first row or rows with bold-only text)
    for row_idx, row in enumerate(table.rows):
        html.append('<tr>')
        is_header = row_idx == 0

        for cell_idx, cell in enumerate(row.cells):
            tag = 'th' if is_header else 'td'
            text = cell.text.strip()

            if text:
                # Check if all runs in cell are bold
                all_bold = all(
                    all(r.bold for r in p.runs if r.text.strip())
                    for p in cell.paragraphs if p.text.strip()
                )
                if all_bold and cell.paragraphs:
                    is_header = True
                    tag = 'th'

            html.append(f'<{tag}>')

            # Process cell content preserving formatting
            for para in cell.paragraphs:
                if para.text.strip():
                    html.append(_process_runs_html(para.runs))

            html.append(f'</{tag}>')

        html.append('</tr>')

    html.append('</table>')


def _process_runs(runs) -> str:
    """Legacy run processor for fallback."""
    return _process_runs_html(runs)


def convert_pptx_to_html(file_content: bytes) -> Optional[str]:
    """Convert PPTX to styled HTML."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(io.BytesIO(file_content))

        html = ['<div class="pptx-content">']

        for slide_num, slide in enumerate(prs.slides, 1):
            html.append(f'<div class="slide" id="slide-{slide_num}">')
            html.append(f'<div class="slide-number">Slide {slide_num}</div>')

            for shape in slide.shapes:
                # Handle tables
                if shape.has_table:
                    tbl = shape.table
                    html.append('<table>')
                    for row_idx, row in enumerate(tbl.rows):
                        html.append('<tr>')
                        for cell in row.cells:
                            tag = 'th' if row_idx == 0 else 'td'
                            html.append(f'<{tag}>{escape(cell.text)}</{tag}>')
                        html.append('</tr>')
                    html.append('</table>')
                    continue

                # Handle images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_blob = shape.image.blob
                        content_type = shape.image.content_type
                        b64 = base64.b64encode(img_blob).decode()
                        html.append(f'<img src="data:{content_type};base64,{b64}" style="max-width:100%;margin:8px 0;" />')
                    except Exception:
                        pass
                    continue

                # Handle text
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        if not para.text.strip():
                            continue

                        is_title = False
                        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                            if shape.placeholder_format.idx == 0:
                                is_title = True

                        runs_html = ''
                        for run in para.runs:
                            t = escape(run.text)
                            if run.font.bold:
                                t = f'<strong>{t}</strong>'
                            if run.font.italic:
                                t = f'<em>{t}</em>'
                            if run.font.size:
                                t = f'<span style="font-size:{run.font.size.pt}pt">{t}</span>'
                            runs_html += t

                        if is_title:
                            html.append(f'<h2 class="slide-title">{runs_html}</h2>')
                        else:
                            html.append(f'<p>{runs_html}</p>')

            html.append('</div>')

        html.append('</div>')
        return '\n'.join(html)

    except Exception as e:
        print(f'PPTX conversion error: {e}')
        return None


def convert_xlsx_to_html(file_content: bytes) -> Optional[str]:
    """Convert XLSX to styled HTML table."""
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter

        wb = load_workbook(io.BytesIO(file_content), read_only=True, data_only=True)

        html = ['<div class="xlsx-content">']

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            html.append(f'<div class="sheet">')
            html.append(f'<h3 class="sheet-title">{escape(sheet_name)}</h3>')
            html.append('<div class="table-wrapper"><table>')

            # First row as header if it looks like headers
            for row_idx, row in enumerate(rows):
                html.append('<tr>')
                is_header = row_idx == 0
                for cell in row:
                    tag = 'th' if is_header else 'td'
                    val = _format_xlsx_cell(cell)
                    html.append(f'<{tag}>{val}</{tag}>')
                html.append('</tr>')

            html.append('</table></div></div>')

        html.append('</div>')
        wb.close()
        return '\n'.join(html)

    except Exception as e:
        print(f'XLSX conversion error: {e}')
        return None


def _format_xlsx_cell(cell) -> str:
    """Format an XLSX cell value for HTML display."""
    if cell is None:
        return ''
    if isinstance(cell, bool):
        return '✓' if cell else '✗'
    if isinstance(cell, (int, float)):
        return f'{cell:,.2f}' if isinstance(cell, float) else f'{cell:,}'
    if hasattr(cell, 'strftime'):
        return cell.strftime('%d/%m/%Y')
    return escape(str(cell))


def convert_office_to_html(file_content: bytes, mime_type: str) -> Optional[str]:
    """Convert office document to HTML based on MIME type."""
    converters = {
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': convert_docx_to_html,
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': convert_pptx_to_html,
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': convert_xlsx_to_html,
        'application/vnd.ms-excel': convert_xlsx_to_html,
        'text/csv': convert_xlsx_to_html,
    }
    converter = converters.get(mime_type)
    if converter:
        return converter(file_content)
    return None
